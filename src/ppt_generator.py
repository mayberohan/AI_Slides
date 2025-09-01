import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE # Added MSO_ANCHOR and MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN # Added PP_ALIGN for text alignment

from src.image_client import fetch_image

def create_presentation(slides_data, output_file, template):
    if not os.path.exists(template):
        raise FileNotFoundError(f"Template not found: {template}")
    
    try:
        prs = Presentation(template)
    except Exception as e:
        print(f"❌ Error loading template: {e}")
        print("Please ensure the template file is not open in another application.")
        return

    # Delete all existing slides from the template
    while len(prs.slides) > 0:
        slide_part = prs.slides._sldIdLst[0]
        prs.part.drop_rel(slide_part.rId)
        del prs.slides._sldIdLst[0]
        
    for slide_data in slides_data:
        title_text = slide_data.get("title", "")
        
        # Try to find a suitable layout - prefer layouts with placeholders
        layout = None
        
        # First, try to find a layout with both content and picture placeholders
        for slide_layout in prs.slide_layouts:
            placeholders = slide_layout.placeholders
            has_title = any(p.placeholder_format.type == PP_PLACEHOLDER.TITLE for p in placeholders)
            has_content = any(p.placeholder_format.type == PP_PLACEHOLDER.BODY for p in placeholders)
            has_picture = any(p.placeholder_format.type == PP_PLACEHOLDER.PICTURE for p in placeholders)
            
            if has_title and has_content and has_picture:
                layout = slide_layout
                print(f"✅ Found layout with title, content, and picture: {slide_layout.name}")
                break
        
        # If no perfect layout, try to find one with title and content
        if not layout:
            for slide_layout in prs.slide_layouts:
                placeholders = slide_layout.placeholders
                has_title = any(p.placeholder_format.type == PP_PLACEHOLDER.TITLE for p in placeholders)
                has_content = any(p.placeholder_format.type == PP_PLACEHOLDER.BODY for p in placeholders)
                
                if has_title and has_content:
                    layout = slide_layout
                    print(f"✅ Found layout with title and content: {slide_layout.name}")
                    break
        
        # Fallback to any layout with a title
        if not layout:
            for slide_layout in prs.slide_layouts:
                placeholders = slide_layout.placeholders
                has_title = any(p.placeholder_format.type == PP_PLACEHOLDER.TITLE for p in placeholders)
                
                if has_title:
                    layout = slide_layout
                    print(f"✅ Found layout with title: {slide_layout.name}")
                    break
        
        # Final fallback
        if not layout:
            layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            print(f"⚠️ Using fallback layout: {layout.name}")

        slide = prs.slides.add_slide(layout)

        # Find placeholders by their type, not by text content
        title_placeholder = None
        body_placeholder = None
        image_placeholder = None
        
        for shape in slide.placeholders:
            placeholder_type = shape.placeholder_format.type
            
            if placeholder_type == PP_PLACEHOLDER.TITLE:
                title_placeholder = shape
            elif placeholder_type == PP_PLACEHOLDER.BODY:
                body_placeholder = shape
            elif placeholder_type == PP_PLACEHOLDER.PICTURE:
                image_placeholder = shape
        
        # Fill the title placeholder
        if title_placeholder and title_placeholder.has_text_frame:
            title_placeholder.text = title_text
            text_frame = title_placeholder.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(36) # Set title font size to 36pt
            print(f"✅ Added title: {title_text} with adjusted font size")
        else:
            print("⚠️ No title placeholder found")
        
        # Add a gap between title and content by adjusting body placeholder position and size
        if title_placeholder and body_placeholder:
            title_bottom = title_placeholder.top + title_placeholder.height
            
            # Set the top of the body placeholder with an additional gap (e.g., 0.2 inches)
            body_placeholder.top = title_bottom + Inches(0.2)
            
            # Adjust the height to ensure it doesn't go off the slide bottom
            # and leaves some margin at the bottom (e.g., 0.2 inches from bottom)
            body_placeholder.height = prs.slide_height - body_placeholder.top - Inches(0.2)

            # Ensure the left and width are reasonable.
            # If there's an image placeholder, adjust width for it.
            if image_placeholder:
                # Assuming image is on the right, content on the left
                body_placeholder.left = Inches(0.5) # Small left margin
                # Calculate width to leave a gap between content and image, and a right margin
                body_placeholder.width = image_placeholder.left - body_placeholder.left - Inches(0.5) 
            else:
                # No image, use full width minus margins
                body_placeholder.left = Inches(0.5)
                body_placeholder.width = prs.slide_width - Inches(1) # 0.5 inch margin on both sides

        # Fill the content placeholder
        if body_placeholder and body_placeholder.has_text_frame:
            text_frame = body_placeholder.text_frame
            text_frame.clear()
            
            # Set text frame properties for left alignment and top anchoring
            text_frame.margin_left = Inches(0.1) # Small left margin
            text_frame.margin_right = Inches(0.1) # Small right margin
            text_frame.vertical_anchor = MSO_ANCHOR.TOP # Ensure content starts from the top
            text_frame.word_wrap = True # Enable word wrap

            content_text = slide_data.get("content", "")
            bullet_points = [p.strip() for p in content_text.split('\n') if p.strip().startswith("-")]

            for i, point in enumerate(bullet_points):
                point = point.lstrip("-").strip()
                if not point: 
                    continue

                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = point
                p.level = 0
                p.font.size = Pt(20) # Increased bullet point font size to 20pt
                p.alignment = PP_ALIGN.LEFT # Set paragraph alignment to left
                # Add some line spacing
                p.space_after = Pt(5) # Add 5pt space after each bullet point
            
            print(f"✅ Added {len(bullet_points)} bullet points with adjusted aesthetics")
        else:
            print("⚠️ No content placeholder found")
        
        # Handle image insertion
        image_path = fetch_image(title_text)
        if image_path and os.path.exists(image_path):
            try:
                if image_placeholder:
                    # Adjust image placeholder position slightly to the left
                    original_left = image_placeholder.left
                    image_placeholder.left = original_left - Inches(0.5) # Shift left by 0.5 inches
                    
                    image_placeholder.insert_picture(image_path)
                    print(f"✅ Image inserted into picture placeholder and shifted left")
                else:
                    # Add image manually to the right side of the slide
                    left = Inches(5.5) # Adjusted left position for manual insertion
                    top = Inches(1.5)
                    width = Inches(3.5)
                    height = Inches(4)
                    
                    # Make sure we don't overlap with existing content
                    if body_placeholder:
                        left = max(left, body_placeholder.left + body_placeholder.width + Inches(0.5))
                    
                    slide.shapes.add_picture(image_path, left, top, width=width, height=height)
                    print(f"✅ Image added manually at position ({left}, {top})")
            except Exception as e:
                print(f"⚠️ Failed to insert image: {e}")
        else:
            print(f"⚠️ No image found for: {title_text}")

        # Add speaker notes
        notes_text = slide_data.get("notes", "")
        if notes_text:
            notes_text_frame = slide.notes_slide.notes_text_frame
            notes_text_frame.text = notes_text
            notes_text_frame.word_wrap = True
            notes_text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            for paragraph in notes_text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(8) # Set notes font size to 8pt
            print(f"✅ Added speaker notes for slide: {title_text} with adjusted aesthetics")

    try:
        prs.save(output_file)
        print(f"✅ Done! Slide deck saved to {output_file}")
    except PermissionError:
        print("❌ Permission denied. Please close the output file if it's open and try again.")
    except Exception as e:
        print(f"❌ An error occurred while saving the presentation: {e}")
