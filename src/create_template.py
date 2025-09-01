import os
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Pt
from PIL import Image, ImageDraw, ImageFont
import io

def create_default_template(filepath: str):
    """
    Creates a PowerPoint template with proper placeholders for title, content, and images.
    """
    if os.path.exists(filepath):
        print(f"Template '{filepath}' already exists.")
        return

    print(f"Creating a new default template at: {filepath}")

    prs = Presentation()
    
    # Use the built-in "Title and Content" layout which has proper placeholders
    # Layout index 1 is typically "Title and Content"
    title_content_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_content_layout)
    
    # The built-in layout already has title and content placeholders
    # We just need to ensure they're properly configured
    
    # Find and configure the title placeholder
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 1:  # Title placeholder
            if shape.has_text_frame:
                shape.text = "Click to add title"
                shape.text_frame.paragraphs[0].font.size = Pt(32)
        elif shape.placeholder_format.type == 2:  # Body/Content placeholder
            if shape.has_text_frame:
                shape.text = "Click to add content"
                shape.text_frame.paragraphs[0].font.size = Pt(18)

    prs.save(filepath)
    print("Template created successfully with proper placeholders.")

def create_advanced_template(filepath: str):
    """
    Creates a more advanced template with title, content, and picture placeholders.
    """
    if os.path.exists(filepath):
        print(f"Advanced template '{filepath}' already exists.")
        return

    print(f"Creating advanced template at: {filepath}")

    prs = Presentation()
    
    # Try to find a layout with picture placeholder (usually "Two Content" or similar)
    layout = None
    for slide_layout in prs.slide_layouts:
        # Look for layouts that have picture placeholders
        has_picture = any(p.placeholder_format.type == 18 for p in slide_layout.placeholders)  # Picture placeholder type
        if has_picture:
            layout = slide_layout
            break
    
    # If no picture layout found, use "Two Content" layout (index 3) or "Title and Content" (index 1)
    if not layout:
        if len(prs.slide_layouts) > 3:
            layout = prs.slide_layouts[3]  # Two Content layout
        else:
            layout = prs.slide_layouts[1]  # Title and Content layout
    
    slide = prs.slides.add_slide(layout)
    
    # Configure placeholders
    for shape in slide.placeholders:
        placeholder_type = shape.placeholder_format.type
        if placeholder_type == 1 and shape.has_text_frame:  # Title
            shape.text = "Click to add title"
            shape.text_frame.paragraphs[0].font.size = Pt(32)
        elif placeholder_type == 2 and shape.has_text_frame:  # Body/Content
            shape.text = "Click to add content"
            shape.text_frame.paragraphs[0].font.size = Pt(18)
        elif placeholder_type == 18:  # Picture placeholder
            # Picture placeholders don't need text configuration
            pass

    prs.save(filepath)
    print("Advanced template created successfully.")

if __name__ == "__main__":
    templates_dir = "templates"
    if not os.path.exists(templates_dir):
        os.makedirs(templates_dir)
    
    create_default_template(os.path.join(templates_dir, "default.pptx"))
