import os
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

def analyze_template(template_path):
    """
    Analyzes a PowerPoint template to show available layouts and placeholders.
    """
    if not os.path.exists(template_path):
        print(f"❌ Template not found: {template_path}")
        return
    
    try:
        prs = Presentation(template_path)
        print(f"📊 Analyzing template: {template_path}")
        print(f"Total layouts: {len(prs.slide_layouts)}")
        print("=" * 50)
        
        for i, layout in enumerate(prs.slide_layouts):
            print(f"\nLayout {i}: {layout.name}")
            print(f"Placeholders: {len(layout.placeholders)}")
            
            for j, placeholder in enumerate(layout.placeholders):
                placeholder_type = placeholder.placeholder_format.type
                type_name = get_placeholder_type_name(placeholder_type)
                
                print(f"  [{j}] Type: {placeholder_type} ({type_name})")
                if placeholder.has_text_frame:
                    print(f"      Has text frame: Yes")
                    if placeholder.text:
                        print(f"      Default text: '{placeholder.text}'")
                else:
                    print(f"      Has text frame: No")
                
                # Show dimensions
                print(f"      Position: ({placeholder.left}, {placeholder.top})")
                print(f"      Size: {placeholder.width} x {placeholder.height}")
        
        print("\n" + "=" * 50)
        print("Recommended layouts for slide generation:")
        
        # Find best layouts
        best_layouts = []
        for i, layout in enumerate(prs.slide_layouts):
            placeholders = layout.placeholders
            has_title = any(p.placeholder_format.type == PP_PLACEHOLDER.TITLE for p in placeholders)
            has_content = any(p.placeholder_format.type == PP_PLACEHOLDER.BODY for p in placeholders)
            has_picture = any(p.placeholder_format.type == PP_PLACEHOLDER.PICTURE for p in placeholders)
            
            score = 0
            features = []
            if has_title:
                score += 3
                features.append("Title")
            if has_content:
                score += 2
                features.append("Content")
            if has_picture:
                score += 1
                features.append("Picture")
            
            if score > 0:
                best_layouts.append((i, layout.name, score, features))
        
        # Sort by score (best first)
        best_layouts.sort(key=lambda x: x[2], reverse=True)
        
        for i, (layout_idx, name, score, features) in enumerate(best_layouts[:3]):
            print(f"{i+1}. Layout {layout_idx}: {name} (Score: {score})")
            print(f"   Features: {', '.join(features)}")
        
    except Exception as e:
        print(f"❌ Error analyzing template: {e}")

def get_placeholder_type_name(placeholder_type):
    """
    Returns a human-readable name for placeholder types.
    """
    type_names = {
        1: "TITLE",
        2: "BODY/CONTENT", 
        3: "CENTER_TITLE",
        4: "SUBTITLE",
        5: "DATE_AND_TIME",
        6: "SLIDE_NUMBER",
        7: "FOOTER",
        8: "HEADER",
        9: "OBJECT",
        10: "CHART",
        11: "TABLE",
        12: "CLIP_ART",
        13: "ORGANIZATION_CHART",
        14: "MEDIA_CLIP",
        15: "SLIDE_IMAGE",
        16: "VERTICAL_OBJECT",
        17: "VERTICAL_TEXT",
        18: "PICTURE"
    }
    return type_names.get(placeholder_type, f"UNKNOWN_{placeholder_type}")

def create_optimized_template(output_path):
    """
    Creates an optimized template specifically for the slide generator.
    """
    print(f"🔧 Creating optimized template: {output_path}")
    
    prs = Presentation()
    
    # Use the "Two Content" layout which typically has title, content, and content placeholders
    # We can use one content placeholder for text and another for images
    if len(prs.slide_layouts) > 3:
        layout = prs.slide_layouts[3]  # Two Content layout
    else:
        layout = prs.slide_layouts[1]  # Title and Content layout
    
    # Add a sample slide to demonstrate the layout
    slide = prs.slides.add_slide(layout)
    
    # Configure the placeholders with sample content
    for shape in slide.placeholders:
        placeholder_type = shape.placeholder_format.type
        if placeholder_type == PP_PLACEHOLDER.TITLE and shape.has_text_frame:
            shape.text = "Sample Title - AI Slide Generator"
        elif placeholder_type == PP_PLACEHOLDER.BODY and shape.has_text_frame:
            shape.text = "• Sample bullet point 1\n• Sample bullet point 2\n• Sample bullet point 3"
    
    prs.save(output_path)
    print(f"✅ Optimized template created: {output_path}")

if __name__ == "__main__":
    # Analyze existing templates
    templates_dir = "templates"
    if os.path.exists(templates_dir):
        for template_file in os.listdir(templates_dir):
            if template_file.endswith('.pptx'):
                template_path = os.path.join(templates_dir, template_file)
                analyze_template(template_path)
                print("\n" + "="*80 + "\n")
    
    # Create an optimized template
    if not os.path.exists(templates_dir):
        os.makedirs(templates_dir)
    
    create_optimized_template(os.path.join(templates_dir, "optimized.pptx"))
